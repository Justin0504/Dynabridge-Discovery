"""PPT Generation module — clones slides from the CozyFit reference PPTX.

Instead of building slides from scratch with python-pptx shapes,
this module copies slides from the original template and replaces
text content with new analysis data. This preserves all formatting,
backgrounds, images, fonts, and layout exactly.

Images are replaced with brand-specific images collected by image_collector.py.
When no collected image is available, the template's original image is kept as-is.

Template slide index map (0-based):
  0  = Cover (Title_Slide)
  1  = Agenda (Blank)
  2  = Approach / Brand Building Process (Text Slide)
  3  = Step 1 Divider (Text Slide)
  4  = Section Header - Capabilities (Overview Slide)
  5  = Content slide - title + 3 bullets + insight + image (Blank)
  13 = Summary slide - title + paragraph + half-image (Text Slide)
  14 = Section Header - Competition (Overview Slide)
  17 = Competitor deep dive - two-column positioning + learnings (Blank)
  23 = Landscape summary - bullets + sidebar text (Blank)
  24 = Summary slide - Competition (Text Slide)
  25 = Section Header - Consumer (Overview Slide)
  91 = Thank You / 谢谢 (Divider Slide)
"""
import copy
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

sys.path.insert(0, str(Path(__file__).parent.parent))
from config import OUTPUT_DIR, PREVIEW_DIR

TEMPLATE_PATH = Path(__file__).parent.parent.parent / "templates" / "cozyfit_reference.pptx"

# Template slide indices (0-based) — which original slide to clone for each type
T_COVER = 0
T_AGENDA = 1
T_APPROACH = 2
T_STEP_DIVIDER = 3
T_SECTION_CAPABILITIES = 4
T_CONTENT = 5          # title + bullets + insight + image
T_CONTENT_ALT = 6      # alternate content layout
T_SUMMARY = 13         # title + summary paragraph + half-image
T_SECTION_COMPETITION = 14
T_COMPETITOR = 17       # two-column: positioning + key learnings
T_LANDSCAPE = 23        # landscape summary (bullets + sidebar)
T_COMP_SUMMARY = 24
T_SECTION_CONSUMER = 25
T_RESEARCH_APPROACH = 26    # Research methodology (label + detail rows)
T_SEGMENT_DIVIDER = 47      # "Market Segmentation" divider
T_SEGMENT_OVERVIEW = 49     # All segments at a glance (names + % + taglines)
T_MEET_SEGMENT = 51         # "Meet the [Segment]" narrative page
T_TARGET_RECOMMENDATION = 76  # "PRIMARY TARGET: [NAME]" with rationale bullets
T_WHY_TARGET = 77           # "WHY [SEGMENT] IS THE RIGHT FOCUS" with bullets
T_ENABLES = 78              # "WHAT THIS CHOICE ENABLES (AND DOES NOT)"
T_CONSUMER_SUMMARY = 79     # Consumer summary (half-text, half-image)
T_FINAL_SUMMARY = 80        # Three-column summary + closing insight
T_THANK_YOU = 91

# Chart slide templates (questionnaire/survey section, slides 27-46)
T_CHART_DIVIDER_DEMO = 27      # "Demographics & Background" divider
T_CHART_SINGLE_HBAR = 34       # Single full-width hbar (e.g., Work Apparel)
T_CHART_DUAL = 35              # Donut left + hbar right (e.g., Purchase Frequency)
T_CHART_SINGLE_VBAR = 30       # Single full-width vbar (e.g., Occupation)
T_CHART_STACKED = 42           # Stacked bar (e.g., Brand Awareness)
T_CHART_DIVIDER_SHOPPING = 33  # "Shopping Habits" divider
T_CHART_DIVIDER_BRAND = 40     # "Brand Evaluation" divider
T_CHART_TABLE = 39             # Challenges table (text-only slide)

# Boilerplate slides
T_SEGMENTATION_INTRO = 48      # "Benefits of segmentation" boilerplate
T_FOCUSING_SEGMENTS = 50       # "FOCUSING ON THE MOST DOMINANT SEGMENTS…"
T_SEGMENT_PROFILE = 52         # Segment respondent profile (demographics layout)
T_CLOSER_LOOK_1 = 53           # "A Closer Look" — premium callout + small icon
T_CLOSER_LOOK_2 = 54           # "A Closer Look" — brand awareness + verbatim quotes
T_CLOSER_LOOK_3 = 56           # "A Closer Look" — 4 lifestyle signal cards
T_CHALLENGES = 55              # Challenges & Pain Points (two tables)
T_SELECTING_TARGET = 75        # "SELECTING [BRAND]'S TARGET AUDIENCE"
T_BRAND_METRICS_DEF = 45       # "Brand Metrics Definitions" boilerplate (GOATClean)

ASSETS_DIR = Path(__file__).parent.parent / "templates" / "assets"


# ── Slide Cloning Engine ─────────────────────────────────────

_src_prs = None


def _get_source():
    """Load the reference PPTX once (cached)."""
    global _src_prs
    if _src_prs is None:
        _src_prs = Presentation(str(TEMPLATE_PATH))
    return _src_prs


def _clone_slide(dst_prs, src_slide_idx):
    """Clone a slide from the reference PPTX into dst_prs.

    Copies all shapes (text boxes, images, groups) and only the
    relationships actually referenced by those shapes. This avoids
    duplicating notesSlide/themeOverride/tags parts that cause ZIP
    corruption.
    """
    src_prs = _get_source()
    src_slide = src_prs.slides[src_slide_idx]

    # Find matching layout in dst by name
    src_layout_name = src_slide.slide_layout.name
    dst_layout = None
    for layout in dst_prs.slide_layouts:
        if layout.name == src_layout_name:
            dst_layout = layout
            break
    if not dst_layout:
        dst_layout = dst_prs.slide_layouts[6]  # Blank fallback

    new_slide = dst_prs.slides.add_slide(dst_layout)

    # Clear auto-generated placeholders from layout
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

    # Step 1: Deep-copy shape elements and collect all rIds they reference
    copied_elements = []
    referenced_rIds = set()
    for shape in src_slide.shapes:
        el = copy.deepcopy(shape._element)
        copied_elements.append(el)
        for attr_el in el.iter():
            for attr_name in list(attr_el.attrib):
                if attr_name == f"{ns_r}id" or attr_name.endswith("}id"):
                    referenced_rIds.add(attr_el.attrib[attr_name])
            if "embed" in attr_el.attrib:
                referenced_rIds.add(attr_el.attrib["embed"])
            if f"{ns_r}embed" in attr_el.attrib:
                referenced_rIds.add(attr_el.attrib[f"{ns_r}embed"])

    # Step 2: Copy relationships — for image/hdphoto parts, create a proper
    # copy in the destination package so partnames don't collide with
    # future add_picture calls (which use next_image_partname).
    SKIP_RELTYPES = {"chart", "tags", "tag", "notesSlide", "themeOverride"}
    IMAGE_RELTYPES = {"image", "hdphoto"}
    rId_map = {}
    for rel in src_slide.part.rels.values():
        if rel.rId not in referenced_rIds:
            continue
        reltype_short = rel.reltype.split("/")[-1]
        if reltype_short in SKIP_RELTYPES:
            continue
        try:
            if reltype_short in IMAGE_RELTYPES and not rel.is_external:
                # Copy image blob into the destination package to avoid
                # cross-package Part references that cause partname collisions
                import io
                src_part = rel._target
                if reltype_short == "image":
                    img_part, new_rId = new_slide.part.get_or_add_image_part(
                        io.BytesIO(src_part.blob)
                    )
                else:
                    # hdphoto — just reference the source directly
                    # (hdphotos don't collide with image indices)
                    new_rId = new_slide.part.rels._add_relationship(
                        rel.reltype, rel._target, rel.is_external
                    )
            else:
                new_rId = new_slide.part.rels._add_relationship(
                    rel.reltype, rel._target, rel.is_external
                )
            rId_map[rel.rId] = new_rId
        except Exception:
            pass

    # Step 3: Remap rIds and remove elements with dangling references
    ns_p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    skip_shapes = []
    for el in copied_elements:
        dangling_elements = []
        for attr_el in el.iter():
            has_dangling = False
            for attr_name in list(attr_el.attrib):
                if attr_name == f"{ns_r}id" or attr_name.endswith("}id"):
                    old_id = attr_el.attrib[attr_name]
                    if old_id in rId_map:
                        attr_el.attrib[attr_name] = rId_map[old_id]
                    elif old_id in referenced_rIds:
                        has_dangling = True
            if "embed" in attr_el.attrib:
                old_id = attr_el.attrib["embed"]
                if old_id in rId_map:
                    attr_el.attrib["embed"] = rId_map[old_id]
                elif old_id in referenced_rIds:
                    has_dangling = True
            if f"{ns_r}embed" in attr_el.attrib:
                old_id = attr_el.attrib[f"{ns_r}embed"]
                if old_id in rId_map:
                    attr_el.attrib[f"{ns_r}embed"] = rId_map[old_id]
                elif old_id in referenced_rIds:
                    has_dangling = True
            if has_dangling:
                dangling_elements.append(attr_el)

        # Remove elements with dangling rIds (chart embeds, tag refs, etc.)
        for dang in dangling_elements:
            parent = dang.getparent()
            if parent is not None:
                parent.remove(dang)

        # Drop the whole top-level shape if it is a graphicFrame whose
        # graphicData no longer has any chart/table/etc. content
        el_tag = el.tag.split("}")[-1]
        if el_tag == "graphicFrame":
            graphic_data = el.find(f".//{ns_a}graphicData")
            if graphic_data is not None and len(graphic_data) == 0:
                skip_shapes.append(el)
                continue

        # Clean up empty custDataLst / extLst containers left behind
        for empty_container_tag in (f"{ns_p}custDataLst", f"{ns_p}extLst", f"{ns_a}extLst"):
            for container in el.findall(f".//{empty_container_tag}"):
                if len(container) == 0:
                    parent = container.getparent()
                    if parent is not None:
                        parent.remove(container)

        new_slide.shapes._spTree.append(el)

    return new_slide


def _find_text_shapes(slide):
    """Return all shapes with text frames, sorted by top then left position."""
    shapes = [s for s in slide.shapes if s.has_text_frame]
    shapes.sort(key=lambda s: (s.top, s.left))
    return shapes


def _set_text_preserve_format(text_frame, new_text):
    """Replace text in a text_frame while preserving per-run and per-paragraph formatting.

    Accepts:
      - str: replaces text line-by-line matching to existing paragraphs/runs
      - list[str]: one string per paragraph, matched 1:1 to existing paragraphs

    Key principle: each original run keeps its font/size/color; we only change .text.
    """
    if isinstance(new_text, list):
        paragraphs = list(text_frame.paragraphs)
        for i, para_text in enumerate(new_text):
            if i < len(paragraphs):
                _replace_para_text(paragraphs[i], para_text)
            else:
                _add_paragraph_after(text_frame, paragraphs[-1], para_text)
        # Remove excess paragraphs
        for i in range(len(new_text), len(paragraphs)):
            p_el = paragraphs[i]._p
            p_el.getparent().remove(p_el)
    else:
        # Split by newlines and match to existing paragraphs/runs
        lines = new_text.split("\n")
        paragraphs = list(text_frame.paragraphs)

        if len(lines) == 1:
            # Single line — distribute across existing runs in first para
            if paragraphs:
                _replace_para_text(paragraphs[0], lines[0])
                for para in paragraphs[1:]:
                    p_el = para._p
                    p_el.getparent().remove(p_el)
        else:
            # Multiple lines — match each line to an existing paragraph
            # If paragraph has multiple runs (different fonts), map lines to runs
            if len(paragraphs) == 1 and len(paragraphs[0].runs) >= len(lines):
                # Single paragraph with multiple runs — map lines to runs
                _replace_runs_text(paragraphs[0], lines)
            else:
                for i, line in enumerate(lines):
                    if i < len(paragraphs):
                        _replace_para_text(paragraphs[i], line)
                    else:
                        _add_paragraph_after(text_frame, paragraphs[-1], line)
                for i in range(len(lines), len(paragraphs)):
                    p_el = paragraphs[i]._p
                    p_el.getparent().remove(p_el)


def _has_cjk(text: str) -> bool:
    """Check if text contains CJK characters."""
    return any('\u4e00' <= c <= '\u9fff' or '\u3400' <= c <= '\u4dbf' for c in text)


CJK_FONT = "Heiti SC"


def _fix_cjk_fonts(prs):
    """Scan all slides and set CJK-compatible font on any run containing CJK text."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if _has_cjk(run.text):
                        _set_cjk_font(run)


def _set_cjk_font(run):
    """Set East Asian font on a run so CJK characters render correctly."""
    from lxml import etree
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }
    rPr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    if rPr is None:
        rPr = etree.SubElement(
            run._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
        run._r.insert(0, rPr)
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        ea = etree.SubElement(
            rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    ea.set('typeface', CJK_FONT)


def _replace_para_text(paragraph, text):
    """Replace text in a paragraph, distributing across existing runs.

    Preserves each run's formatting (font, size, color, bold).
    Sets CJK-compatible font when text contains Chinese characters.
    """
    runs = paragraph.runs
    if not runs:
        paragraph.text = text
        return

    if len(runs) == 1:
        runs[0].text = text
        if _has_cjk(text):
            _set_cjk_font(runs[0])
    else:
        runs[0].text = text
        if _has_cjk(text):
            _set_cjk_font(runs[0])
        for run in runs[1:]:
            run.text = ""


def _replace_runs_text(paragraph, texts):
    """Replace text run-by-run, one text per run, preserving each run's formatting."""
    runs = paragraph.runs
    for i, text in enumerate(texts):
        if i < len(runs):
            runs[i].text = text
    # Clear excess runs
    for i in range(len(texts), len(runs)):
        runs[i].text = ""


def _add_paragraph_after(text_frame, template_para, text):
    """Add a new paragraph after template_para, copying its formatting."""
    new_p = copy.deepcopy(template_para._p)
    template_para._p.addnext(new_p)
    from pptx.text.text import _Paragraph
    para = _Paragraph(new_p, template_para._parent)
    _replace_para_text(para, text)


# ── Text Truncation ──────────────────────────────────────────

def _truncate(text, max_chars):
    """Truncate text to max_chars, preferring sentence boundaries."""
    if not text or len(text) <= max_chars:
        return text
    # Prefer cutting at last sentence end (period) before limit
    last_period = text[:max_chars].rfind(".")
    if last_period > max_chars * 0.4:
        return text[:last_period + 1]
    # Fall back to last comma or semicolon
    for sep in (",", ";", " —", " –"):
        pos = text[:max_chars].rfind(sep)
        if pos > max_chars * 0.4:
            return text[:pos] + "."
    # Last resort: word boundary
    cut = text[:max_chars].rfind(" ")
    if cut < max_chars // 2:
        cut = max_chars
    return text[:cut].rstrip(" ,;:") + "."


# ── Image Replacement ───────────────────────────────────────

def _replace_slide_image(slide, image_path: Path, replace_background=False):
    """Replace a picture shape on a slide with a new image.

    By default replaces the first non-background picture. With
    replace_background=True, replaces the largest (background) picture
    instead — used for "Meet the Segment" hero slides.

    Pre-crops the image file with PIL to exactly match the box aspect
    ratio (cover + top-bias), then inserts the cropped image at the
    exact box dimensions.
    """
    if not image_path or not Path(image_path).exists():
        return

    from PIL import Image
    from pptx.shapes.picture import Picture

    SLIDE_AREA = 12192000 * 6858000  # 16:9 widescreen

    # Collect all picture shapes first to avoid mutation-during-iteration
    pictures = []
    for shape in slide.shapes:
        if isinstance(shape, Picture):
            box_w, box_h = shape.width, shape.height
            is_bg = (box_w * box_h) / SLIDE_AREA > 0.9
            if replace_background and not is_bg:
                continue
            if not replace_background and is_bg:
                continue
            pictures.append(shape)

    if not pictures:
        return

    shape = pictures[0]
    box_left, box_top = shape.left, shape.top
    box_w, box_h = shape.width, shape.height
    if box_h == 0:
        return
    box_ratio = box_w / box_h

    try:
        img = Image.open(str(image_path))
        if img.mode in ("RGBA", "P", "LA"):
            img = img.convert("RGBA")
        else:
            img = img.convert("RGB")
        img_w, img_h = img.size
    except Exception:
        return

    img_ratio = img_w / img_h

    # Crop image to match box aspect ratio (cover mode)
    if abs(img_ratio - box_ratio) > 0.05:
        if img_ratio > box_ratio:
            new_w = int(img_h * box_ratio)
            offset = (img_w - new_w) // 2
            img = img.crop((offset, 0, offset + new_w, img_h))
        else:
            new_h = int(img_w / box_ratio)
            top_offset = int((img_h - new_h) * 0.15)
            img = img.crop((0, top_offset, img_w, top_offset + new_h))

    cropped_path = image_path.parent / f"_cropped_{image_path.stem}.png"
    img.save(str(cropped_path), format="PNG")
    img.close()

    # Get a properly-registered ImagePart for the cropped image
    image_part, rId = slide.part.get_or_add_image_part(str(cropped_path))

    # Update the existing Picture's blip to reference the new image,
    # rather than remove+add which risks partname collisions with
    # source Parts from _clone_slide that share the same package
    ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    ns_a_uri = "http://schemas.openxmlformats.org/drawingml/2006/main"
    blip = shape._element.find(f".//{{{ns_a_uri}}}blip")
    if blip is not None:
        blip.set(f"{ns_r}embed", rId)
        # Remove artistic effect layers that reference old image parts
        for child in list(blip):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "imgLayer":
                blip.remove(child)
        return

    # Fallback: remove and recreate
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(str(cropped_path), box_left, box_top, box_w, box_h)
    return


def _replace_card_images(slide, img_pool):
    """Replace multiple card-sized images on a slide (e.g., Closer Look 3).

    Finds all Picture shapes between 5% and 20% of slide area (the 4 lifestyle
    cards) and replaces each with a different brand image.
    """
    from PIL import Image
    from pptx.shapes.picture import Picture

    SLIDE_AREA = 12192000 * 6858000

    cards = []
    for shape in slide.shapes:
        if isinstance(shape, Picture):
            area_pct = (shape.width * shape.height) / SLIDE_AREA
            if 0.05 < area_pct < 0.20:
                cards.append(shape)

    for shape in cards:
        image_path = img_pool.next_brand()
        if not image_path or not Path(image_path).exists():
            continue

        box_w, box_h = shape.width, shape.height
        if box_h == 0:
            continue
        box_ratio = box_w / box_h

        try:
            img = Image.open(str(image_path))
            if img.mode in ("RGBA", "P", "LA"):
                img = img.convert("RGBA")
            else:
                img = img.convert("RGB")
            img_w, img_h = img.size
        except Exception:
            continue

        img_ratio = img_w / img_h
        if abs(img_ratio - box_ratio) > 0.05:
            if img_ratio > box_ratio:
                new_w = int(img_h * box_ratio)
                offset = (img_w - new_w) // 2
                img = img.crop((offset, 0, offset + new_w, img_h))
            else:
                new_h = int(img_w / box_ratio)
                top_offset = int((img_h - new_h) * 0.15)
                img = img.crop((0, top_offset, img_w, top_offset + new_h))

        cropped_path = image_path.parent / f"_cropped_{image_path.stem}.png"
        img.save(str(cropped_path), format="PNG")
        img.close()

        image_part, rId = slide.part.get_or_add_image_part(str(cropped_path))
        ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
        ns_a_uri = "http://schemas.openxmlformats.org/drawingml/2006/main"
        blip = shape._element.find(f".//{{{ns_a_uri}}}blip")
        if blip is not None:
            blip.set(f"{ns_r}embed", rId)
            for child in list(blip):
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if tag == "imgLayer":
                    blip.remove(child)


class _ImagePool:
    """Manages a pool of collected images, handing them out one at a time.

    Images are categorized: brand images are used first for content slides,
    product images for competitor slides, lifestyle for summaries.
    """

    def __init__(self, images: dict = None):
        self._images = images or {}
        self._brand_idx = 0
        self._product_idx = 0
        self._lifestyle_idx = 0
        self._all_idx = 0

    def next_brand(self) -> Path | None:
        """Get next brand image, cycling through available images."""
        imgs = self._images.get("brand", [])
        if not imgs:
            return self.next_any()
        img = imgs[self._brand_idx % len(imgs)]
        self._brand_idx += 1
        return img

    def next_product(self) -> Path | None:
        """Get next product image."""
        imgs = self._images.get("product", [])
        if not imgs:
            return self.next_brand()
        img = imgs[self._product_idx % len(imgs)]
        self._product_idx += 1
        return img

    def next_lifestyle(self) -> Path | None:
        """Get next lifestyle/stock image."""
        imgs = self._images.get("lifestyle", [])
        if not imgs:
            return self.next_brand()
        img = imgs[self._lifestyle_idx % len(imgs)]
        self._lifestyle_idx += 1
        return img

    def next_any(self) -> Path | None:
        """Get any available image."""
        imgs = self._images.get("all", [])
        if not imgs:
            return None
        img = imgs[self._all_idx % len(imgs)]
        self._all_idx += 1
        return img

    def has_images(self) -> bool:
        return bool(self._images.get("all"))


# ── High-level Slide Builders ────────────────────────────────

def _build_cover(prs, brand_name, date_str):
    """Clone cover slide, replace brand name and date.

    Original cover has one text frame with 2 runs in 1 paragraph:
      Run 0: "CozyFit"  (Montserrat 60pt)
      Run 1: "Brand Discovery" (Montserrat 35pt, preceded by line break)
    We replace run-by-run to preserve each font size.
    """
    slide = _clone_slide(prs, T_COVER)
    shapes = _find_text_shapes(slide)
    if len(shapes) >= 1:
        tf = shapes[0].text_frame
        para = tf.paragraphs[0]
        runs = para.runs
        if len(runs) >= 2:
            # Run 0 = brand name, Run 1 = subtitle
            runs[0].text = brand_name
            runs[1].text = "\nBrand Discovery"
        else:
            _set_text_preserve_format(tf, f"{brand_name}\nBrand Discovery")
    if len(shapes) >= 2:
        _set_text_preserve_format(shapes[1].text_frame, date_str)
    return slide


def _build_agenda(prs):
    """Clone agenda slide (no text changes needed — it's generic)."""
    return _clone_slide(prs, T_AGENDA)


def _build_approach(prs):
    """Clone the 'Our Brand Building Process' approach slide."""
    return _clone_slide(prs, T_APPROACH)


def _build_step_divider(prs):
    """Clone the 'Step 1 – Discovery' divider."""
    return _clone_slide(prs, T_STEP_DIVIDER)


def _build_section_header(prs, section_type):
    """Clone a section header. section_type: 'capabilities'|'competition'|'consumer'."""
    idx_map = {
        "capabilities": T_SECTION_CAPABILITIES,
        "competition": T_SECTION_COMPETITION,
        "consumer": T_SECTION_CONSUMER,
    }
    return _clone_slide(prs, idx_map.get(section_type, T_SECTION_CAPABILITIES))


def _build_content_slide(prs, title, bullets, insight_text, template_idx=T_CONTENT):
    """Clone a content slide (title + bullets + insight + image).

    Template shape layout (sorted by position):
      Shape 0 (top): Title — ALL CAPS, orange, Montserrat Bold
      Shape 1 (middle): Bullets — 3 paragraphs, Montserrat, space_before/after
      Shape 2 (bottom): Insight — teal/blue text, single paragraph
      Shape 3: Image (preserved as-is from template)

    Character limits (from original CozyFit template):
      Title: ~55 chars, Bullets: ~100 chars each, Insight: ~90 chars
    """
    slide = _clone_slide(prs, template_idx)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, _truncate(title, 55))
    if len(shapes) >= 2:
        if isinstance(bullets, list):
            bullets = [_truncate(b, 85) for b in bullets[:3]]
        else:
            bullets = [_truncate(bullets, 85)]
        _set_text_preserve_format(shapes[1].text_frame, bullets)
    if len(shapes) >= 3:
        _set_text_preserve_format(shapes[2].text_frame, _truncate(insight_text, 85))

    return slide


def _build_competitor_slide(prs, name, positioning_bullets, learnings_bullets):
    """Clone a competitor deep-dive slide (two-column: positioning + learnings).

    Template shape layout:
      Shape 0: Title — "DICKIES — POSITIONING & KEY LEARNINGS"
      Shape 1: Left column — "POSITIONING\nbullet1\nbullet2\n..."
      Shape 2: Right column — "KEY LEARNINGS\nbullet1\nbullet2\n..."
      Shape 3+: Images (preserved)
    """
    slide = _clone_slide(prs, T_COMPETITOR)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, _truncate(f"{name.upper()} — POSITIONING & KEY LEARNINGS", 60))

    if len(shapes) >= 2:
        positioning_text = ["POSITIONING"] + [_truncate(b, 90) for b in positioning_bullets[:3]]
        _set_text_preserve_format(shapes[1].text_frame, positioning_text)

    if len(shapes) >= 3:
        learnings_text = ["KEY LEARNINGS"] + [_truncate(b, 90) for b in learnings_bullets[:3]]
        _set_text_preserve_format(shapes[2].text_frame, learnings_text)

    return slide


def _build_landscape_slide(prs, title, bullets, sidebar_text):
    """Clone the landscape summary slide (slide 24 pattern)."""
    slide = _clone_slide(prs, T_LANDSCAPE)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, _truncate(title, 60))
    if len(shapes) >= 2:
        if isinstance(bullets, list):
            bullets = [_truncate(b, 100) for b in bullets]
        else:
            bullets = [_truncate(bullets, 100)]
        _set_text_preserve_format(shapes[1].text_frame, bullets)
    if len(shapes) >= 3:
        _set_text_preserve_format(shapes[2].text_frame, _truncate(sidebar_text, 180))
        # Reduce font size for sidebar callout to prevent overflow
        for para in shapes[2].text_frame.paragraphs:
            for run in para.runs:
                if run.font.size and run.font.size > Pt(16):
                    run.font.size = Pt(14)

    return slide


def _build_summary_slide(prs, title, summary_text, template_idx=T_SUMMARY):
    """Clone a summary slide (title + flowing paragraph + half-image)."""
    slide = _clone_slide(prs, template_idx)
    shapes = _find_text_shapes(slide)

    # Summary slide has 2 text shapes: paragraph body and title
    # They may be in different order depending on position sort
    title_shape = None
    body_shape = None
    for s in shapes:
        text = s.text_frame.text.strip().upper()
        if "SUMMARY" in text or len(text) < 40:
            title_shape = s
        else:
            body_shape = s

    if title_shape:
        _set_text_preserve_format(title_shape.text_frame, _truncate(title, 40))
    if body_shape:
        _set_text_preserve_format(body_shape.text_frame, _truncate(summary_text, 260))

    return slide


def _build_thank_you(prs):
    """Clone the Thank You slide."""
    return _clone_slide(prs, T_THANK_YOU)


# ── Consumer Slide Builders ─────────────────────────────────

def _build_research_approach(prs, research_items):
    """Clone research approach slide (slide 26 pattern).

    Template has label+detail rows: Format, Participants, Analysis, Timing.
    Each row is two text shapes side by side (label left, detail right).
    """
    slide = _clone_slide(prs, T_RESEARCH_APPROACH)
    shapes = _find_text_shapes(slide)

    # Shape 0 = title, then pairs of (label, detail)
    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "RESEARCH APPROACH")

    # Map research_items to the label+detail shape pairs
    pair_idx = 0
    for item in research_items[:5]:
        label_shape_idx = 2 + pair_idx * 2
        detail_shape_idx = label_shape_idx - 1
        # Template order: detail shape comes before label in position sort
        # Actual layout: shapes alternate detail(left-wide) and label(left-narrow)
        if detail_shape_idx < len(shapes) and label_shape_idx < len(shapes):
            _set_text_preserve_format(shapes[label_shape_idx].text_frame, item.get("label", ""))
            _set_text_preserve_format(shapes[detail_shape_idx].text_frame, _truncate(item.get("detail", ""), 200))
        pair_idx += 1

    return slide


def _build_segment_overview(prs, segments):
    """Clone segment overview slide (slide 49 pattern).

    Shows all segments at a glance: name, %, tagline for each.
    Template has 5 columns with: percentage text, image, name, tagline.
    """
    slide = _clone_slide(prs, T_SEGMENT_OVERVIEW)
    shapes = _find_text_shapes(slide)

    # Shape 0 = title
    if shapes:
        _set_text_preserve_format(shapes[0].text_frame, "CONSUMER SEGMENTS AT A GLANCE")

    # Shapes 1-5 = percentages, 11-15 = names, 16-20 = taglines
    # Find percentage shapes (short text, typically "27%")
    pct_shapes = []
    name_shapes = []
    tagline_shapes = []
    for s in shapes[1:]:
        text = s.text_frame.text.strip()
        if text.endswith("%") and len(text) <= 4:
            pct_shapes.append(s)
        elif len(text) < 30 and not text.endswith("%"):
            name_shapes.append(s)
        elif len(text) > 30:
            tagline_shapes.append(s)

    for i, seg in enumerate(segments[:5]):
        if i < len(pct_shapes):
            _set_text_preserve_format(pct_shapes[i].text_frame, f"{seg.get('size_pct', '?')}%")
        if i < len(name_shapes):
            _set_text_preserve_format(name_shapes[i].text_frame, seg.get("name", f"Segment {i+1}"))
        if i < len(tagline_shapes):
            _set_text_preserve_format(tagline_shapes[i].text_frame, _truncate(seg.get("tagline", ""), 80))

    return slide


def _build_meet_segment(prs, segment):
    """Clone 'Meet the [Segment]' slide (slide 51 pattern).

    Full-bleed background image with overlay text:
      Shape 0: background (skip)
      Shape 1: background image
      Shape 2: Segment name (ALL CAPS, large)
      Shape 3: Tagline (one line)
      Shape 4: Narrative paragraph (5-7 sentences)
    """
    slide = _clone_slide(prs, T_MEET_SEGMENT)
    shapes = _find_text_shapes(slide)

    name = segment.get("name", "SEGMENT")
    tagline = segment.get("tagline", "")
    narrative = segment.get("narrative", "")

    # Find the shapes by content length pattern
    for s in shapes:
        text = s.text_frame.text.strip()
        if not text:
            continue
        if text.isupper() and len(text) < 30:
            _set_text_preserve_format(s.text_frame, name.upper())
        elif len(text) < 80 and not text.isupper() and "Meet" not in text:
            _set_text_preserve_format(s.text_frame, _truncate(tagline, 70))
        elif len(text) > 80 or "Meet" in text:
            _set_text_preserve_format(s.text_frame, _truncate(narrative, 500))

    return slide


def _build_segment_closer_look(prs, segment, slide_num=1):
    """Build 'A Closer Look' slide using the real template layouts.

    Slide_num determines template and data:
      1: T_CLOSER_LOOK_1 (slide 53) — premium/driver callout text
      2: T_CLOSER_LOOK_2 (slide 54) — brand awareness + verbatim quotes
      3: T_CLOSER_LOOK_3 (slide 56) — 4 lifestyle signal cards
    """
    name = segment.get("name", "SEGMENT")
    mini = segment.get("mini_tables", {})
    lifestyle = segment.get("lifestyle_signals", [])
    title = f"{name.upper()} – A CLOSER LOOK"

    if slide_num == 1:
        slide = _clone_slide(prs, T_CLOSER_LOOK_1)
        shapes = _find_text_shapes(slide)
        drivers = mini.get("purchase_drivers", [])
        motivations = segment.get("key_motivations", [])
        unmet = segment.get("unmet_needs", "")
        description = segment.get("description", "")
        wtp = segment.get("willingness_to_pay", "")
        strategic = segment.get("strategic_value", "")

        # Build callout from whatever data is available
        callout = ""
        if drivers:
            top = " and ".join(f"{d['item'].lower()} ({d['pct']}%)" for d in drivers[:2])
            callout = f"Top purchase drivers are {top}. {drivers[0]['item']} ({drivers[0]['pct']}%) leads all drivers."
        elif motivations:
            items = ", ".join(str(m).lower() for m in motivations[:3])
            callout = f"Key motivations: {items}."
            if unmet:
                callout += f" Unmet need: {unmet}"
        elif description:
            callout = description

        stat = ""
        pain = mini.get("pain_points", [])
        if pain:
            stat = f"{pain[0]['item'].lower()} ({pain[0]['pct']}%) is the top pain point"
        elif wtp:
            stat = f"Willingness to pay: {wtp}"
        elif strategic:
            stat = strategic[:80]

        # Assign text by shape role: title gets segment name, the two largest
        # remaining text areas get callout and stat, everything else is cleared
        # to prevent overlap from the many small template text boxes.
        title_done = False
        candidates = []
        for s in shapes:
            text = s.text_frame.text.strip()
            if not title_done and ("CLOSER LOOK" in text.upper() or "ENDURANCE" in text.upper()
                                   or "PROFILE" in text.upper()):
                _set_text_preserve_format(s.text_frame, _truncate(title, 55))
                title_done = True
            elif "base" in text.lower():
                continue  # leave sample-size footer as-is
            elif text:
                area = s.width * s.height
                candidates.append((area, s))

        # Sort by area descending — largest gets callout, second gets stat
        candidates.sort(key=lambda x: x[0], reverse=True)
        for i, (area, s) in enumerate(candidates):
            if i == 0 and callout:
                # Estimate max chars from box width (roughly 10 chars per inch at 16pt)
                max_chars = max(40, int(s.width / 914400 * 10))
                _set_text_preserve_format(s.text_frame, _truncate(callout, max_chars))
            elif i == 1 and stat:
                max_chars = max(20, int(s.width / 914400 * 10))
                _set_text_preserve_format(s.text_frame, _truncate(stat, max_chars))
            else:
                # Clear remaining small boxes to prevent clutter/overlap
                for p in s.text_frame.paragraphs:
                    for r in p.runs:
                        r.text = ""
        return slide

    elif slide_num == 2:
        slide = _clone_slide(prs, T_CLOSER_LOOK_2)
        shapes = _find_text_shapes(slide)
        shop = segment.get("shopping_behavior", {})
        needs = segment.get("top_needs", [])
        pain_points = segment.get("pain_points", [])

        # Insert a visual in the area where the template's chart used to be
        # (graphicFrame was stripped during clone). Prefer real percentage data
        # from mini_tables; fall back to ranked key_motivations/unmet_needs.
        categories, values = [], []
        for src_key in ("purchase_drivers", "pain_points", "pre_purchase"):
            data = mini.get(src_key) or []
            if isinstance(data, list) and data and isinstance(data[0], dict) and "pct" in data[0]:
                for item in data[:5]:
                    label = str(item.get("item", ""))[:24]
                    pct = item.get("pct", 0)
                    if label and pct:
                        categories.append(label)
                        values.append(float(pct))
                break

        if not categories:
            motivations = segment.get("key_motivations") or []
            if isinstance(motivations, list) and motivations:
                ranked = [(str(m)[:24], 100 - i * 20) for i, m in enumerate(motivations[:5])]
                categories = [r[0] for r in ranked]
                values = [r[1] for r in ranked]

        if categories and values:
            from pipeline.chart_renderer import render_hbar
            import tempfile
            tmp_chart = Path(tempfile.mkstemp(suffix="_closer_bar.png")[1])
            # Narrower than original graphicFrame (5.27M EMU) so it doesn't
            # collide with the quote bubble at x=4583380
            render_hbar(categories, values, output_path=tmp_chart, size=(850, 1000))
            slide.shapes.add_picture(
                str(tmp_chart),
                Emu(426582), Emu(1606658),
                width=Emu(4056798), height=Emu(4795527),
            )

        # Build richer fallback text from segment-level data
        channels = segment.get("channels", [])
        touchpoints = segment.get("media_touchpoints", [])
        wtp = segment.get("willingness_to_pay", "")
        description = segment.get("description", "")
        unmet = segment.get("unmet_needs", "")

        channel_detail = ""
        if channels:
            channel_detail = f"Channels: {', '.join(channels[:3])}"
        if touchpoints:
            channel_detail += f" | Media: {', '.join(touchpoints[:3])}"
        if wtp:
            channel_detail += f" | WTP: {wtp}"

        # Map text shapes by position and content
        for s in shapes:
            text = s.text_frame.text.strip()
            if "CLOSER LOOK" in text.upper() or "ENDURANCE" in text.upper():
                _set_text_preserve_format(s.text_frame, _truncate(title, 55))
            elif "Brand Awareness" in text:
                label = "Key Motivations" if segment.get("key_motivations") else "Shopping Behavior"
                _set_text_preserve_format(s.text_frame, label)
            elif "Top 3" in text or "important features" in text.lower():
                if needs:
                    label = f"Top needs: {', '.join(needs[:3])}"
                elif channels:
                    label = f"Key channels: {', '.join(channels[:3])}"
                else:
                    label = description[:80] if description else name
                _set_text_preserve_format(s.text_frame, _truncate(label, 80))
            elif "comfort" in text.lower() or "All-day" in text:
                pre = mini.get("pre_purchase", [])
                detail = " | ".join(f"{p['item']} ({p['pct']}%)" for p in pre[:3]) if pre else ""
                if shop:
                    detail = f"Channel: {shop.get('primary_channel', 'N/A')} | Spend: {shop.get('annual_spend', 'N/A')} | {detail}"
                if not detail:
                    detail = channel_detail or description[:150]
                _set_text_preserve_format(s.text_frame, _truncate(detail, 150))
            elif "anything else" in text.lower() or "experience" in text.lower():
                prompt = unmet[:80] if unmet else "What challenges or pain points do they face?"
                _set_text_preserve_format(s.text_frame, _truncate(prompt, 80))
            elif text.startswith('"') or "stress" in text.lower() or "wish" in text.lower():
                # Verbatim quote bubbles — fill with pain points or motivations/needs
                quotes = list(pain_points) if pain_points else []
                if not quotes:
                    for m in segment.get("key_motivations", []):
                        quotes.append(f"I need {str(m).lower()}")
                    if unmet:
                        for sentence in unmet.split(". "):
                            if sentence.strip():
                                quotes.append(sentence.strip())
                idx = 0
                for s2 in shapes:
                    t2 = s2.text_frame.text.strip()
                    if t2.startswith('"') or (len(t2) > 20 and s2.top > 3000000):
                        if idx < len(quotes):
                            _set_text_preserve_format(s2.text_frame, _truncate(f'"{quotes[idx]}"', 100))
                        else:
                            _set_text_preserve_format(s2.text_frame, "")
                        idx += 1
                break
        return slide

    else:
        slide = _clone_slide(prs, T_CLOSER_LOOK_3)
        shapes = _find_text_shapes(slide)
        # Slide 56 layout: title + 4 lifestyle text boxes at bottom

        # If no lifestyle_signals, synthesize from channels/motivations/touchpoints
        if not lifestyle or all(not l.get("detail") for l in lifestyle):
            lifestyle = []
            channels = segment.get("channels", [])
            touchpoints = segment.get("media_touchpoints", [])
            motivations = segment.get("key_motivations", [])
            for ch in channels[:2]:
                lifestyle.append({"category": "Channel", "detail": str(ch)})
            for tp in touchpoints[:1]:
                lifestyle.append({"category": "Media", "detail": str(tp)})
            for m in motivations[:1]:
                lifestyle.append({"category": "Motivation", "detail": str(m)})

        while len(lifestyle) < 4:
            lifestyle.append({"category": "", "detail": ""})

        for s in shapes:
            text = s.text_frame.text.strip()
            if "CLOSER LOOK" in text.upper() or "ENDURANCE" in text.upper():
                _set_text_preserve_format(s.text_frame, _truncate(title, 55))

        # Find the 4 bottom text boxes (top > 4500000, sorted by left)
        bottom_shapes = sorted(
            [s for s in shapes if s.top > 4500000 and s.text_frame.text.strip() and "base" not in s.text_frame.text.lower()],
            key=lambda s: s.left,
        )
        for idx, s in enumerate(bottom_shapes[:4]):
            if idx < len(lifestyle) and lifestyle[idx].get("detail"):
                _set_text_preserve_format(s.text_frame, _truncate(lifestyle[idx]["detail"], 90))
            else:
                _set_text_preserve_format(s.text_frame, "")

        return slide


def _build_segment_profile(prs, segment):
    """Clone the respondent profile slide for a segment.

    Template idx 52 has 2 chart graphicFrames (stripped during clone),
    demographic icons, and label text boxes.  We replace the stripped
    charts with a rendered demographics bar, and update the labels
    with whatever structured data we can extract.
    """
    slide = _clone_slide(prs, T_SEGMENT_PROFILE)
    name = segment.get("name", "SEGMENT")
    demo_raw = segment.get("demographics", "")
    size_pct = segment.get("size_pct", "")
    channels = segment.get("channels", [])
    wtp = segment.get("willingness_to_pay", "")

    shapes = _find_text_shapes(slide)
    for s in shapes:
        text = s.text_frame.text.strip()
        if "RESPONDENT PROFILE" in text.upper() or "ENDURANCE" in text.upper():
            _set_text_preserve_format(s.text_frame, f"{name.upper()} – RESPONDENT PROFILE")
        elif "base" in text.lower() and "n" in text.lower():
            _set_text_preserve_format(s.text_frame, f"Segment size: {size_pct}% of audience" if size_pct else "")

    # Parse simple tokens from demographics string
    demo_lower = demo_raw.lower() if demo_raw else ""
    demo_items = {
        "Generation": "",
        "Marital Status": "",
        "Household Income": "",
        "Race / Ethnicity": "",
    }
    import re as _re
    age_match = _re.search(r'(\d{1,2}[\s–\-]+\d{1,2})', demo_raw)
    if age_match:
        demo_items["Generation"] = f"Age {age_match.group(1)}"
    income_match = _re.search(r'(\$[\d,]+[kK]?\s*[\-–]\s*\$[\d,]+[kK]?)', demo_raw)
    if income_match:
        demo_items["Household Income"] = income_match.group(1)
    for kw in ("female", "male", "women", "men"):
        if kw in demo_lower:
            demo_items["Race / Ethnicity"] = kw.capitalize()
            break

    # Update matching label shapes
    for s in shapes:
        text = s.text_frame.text.strip()
        for label, value in demo_items.items():
            if text == label:
                continue  # keep the header labels as-is
        # Replace percentage placeholders with parsed values or dashes
        if _re.match(r'^\d{1,3}%$', text):
            _set_text_preserve_format(s.text_frame, "—")
        elif text in ("Married or domestic partnership", "Single, never married",
                       "Widowed, divorced or separated"):
            pass  # keep as template labels
        elif text.startswith("Low") or text.startswith("High") or text.startswith("Upper"):
            pass  # keep income bracket labels

    # Insert a demographics summary bar chart where charts were stripped
    categories = []
    values = []
    if channels:
        for i, ch in enumerate(channels[:5]):
            categories.append(str(ch)[:20])
            values.append(100 - i * 15)
    if categories and values:
        from pipeline.chart_renderer import render_hbar
        import tempfile
        tmp_chart = Path(tempfile.mkstemp(suffix="_profile_bar.png")[1])
        render_hbar(categories, values, output_path=tmp_chart,
                    question="Top Channels", size=(1330, 1000))
        slide.shapes.add_picture(
            str(tmp_chart),
            Emu(393229), Emu(1483388),
            width=Emu(3551494), height=Emu(2674705),
        )

    return slide


def _build_segment_challenges(prs, segment):
    """Clone the challenges & pain points slide for a segment.

    Template idx 55 has two tables:
      - Shape 1: 16×1 verbatim quote table
      - Shape 3: 8×2 pain-point + percentage table
    We fill them from unmet_needs, key_motivations, and description.
    """
    slide = _clone_slide(prs, T_CHALLENGES)
    name = segment.get("name", "SEGMENT")
    unmet = segment.get("unmet_needs", "")
    motivations = segment.get("key_motivations", [])
    description = segment.get("description", "")

    shapes = _find_text_shapes(slide)
    for s in shapes:
        text = s.text_frame.text.strip()
        if "CHALLENGES" in text.upper() or "ENDURANCE" in text.upper():
            _set_text_preserve_format(s.text_frame, f"{name.upper()} – KEY NEEDS & CHALLENGES")
        elif text == "Pain Points":
            _set_text_preserve_format(s.text_frame, "Top Needs")

    # Build quote lines from unmet_needs
    quote_lines = []
    if unmet:
        for sentence in unmet.replace(". ", ".\n").split("\n"):
            sentence = sentence.strip()
            if sentence:
                quote_lines.append(f'"{sentence}"')
    while len(quote_lines) < 4 and motivations:
        m = motivations.pop(0)
        quote_lines.append(f'"{m}"')

    # Build needs rows from key_motivations + parsed unmet_needs
    need_rows = []
    for m in segment.get("key_motivations", []):
        need_rows.append((str(m)[:50], ""))
    if not need_rows and unmet:
        for part in unmet.split(","):
            part = part.strip()
            if part:
                need_rows.append((part[:50], ""))

    # Fill tables
    for sh in slide.shapes:
        if hasattr(sh, "has_table") and sh.has_table:
            t = sh.table
            cols = len(t.columns)
            rows = len(t.rows)
            if cols == 1 and rows > 4:
                # Verbatim quotes table
                for ri in range(rows):
                    cell = t.cell(ri, 0)
                    if ri == 0:
                        cell.text = "What are the unmet needs of this segment?"
                    elif ri - 1 < len(quote_lines):
                        cell.text = quote_lines[ri - 1]
                    else:
                        cell.text = ""
            elif cols == 2:
                # Pain points / needs table
                for ri in range(rows):
                    if ri < len(need_rows):
                        t.cell(ri, 0).text = need_rows[ri][0]
                        t.cell(ri, 1).text = need_rows[ri][1]
                    else:
                        t.cell(ri, 0).text = ""
                        t.cell(ri, 1).text = ""

    return slide


def _build_brand_metrics_def(prs):
    """Clone brand metrics definitions boilerplate slide."""
    slide = _clone_slide(prs, T_CHART_TABLE)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "BRAND METRICS DEFINITIONS")

    definitions = (
        "Aided Awareness: % of consumers who recognize your brand when prompted.\n"
        "Purchase: % of consumers who have bought your brand's products within a specific timeframe.\n"
        "Satisfaction: % of purchasers who report being satisfied with the product.\n"
        "Recommendation: % of purchasers likely to recommend the brand to others."
    )

    _remove_chart_shapes(slide, clean_region=True)
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(
        Emu(419100), Emu(1400000), Emu(11353800), Emu(4000000)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in definitions.split("\n"):
        p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x29, 0x25, 0x24)

    return slide


def _build_segmentation_intro(prs):
    """Clone segmentation benefits boilerplate slide."""
    slide = _clone_slide(prs, T_SEGMENTATION_INTRO)
    return slide


def _build_focusing_segments(prs, segments):
    """Clone 'FOCUSING ON THE MOST DOMINANT SEGMENTS…' slide (slide 50).

    Shows segments with percentages and taglines, highlighting the dominant ones.
    Template has: title + 5 percentage auto-shapes + 5 name text boxes + 5 taglines.
    """
    slide = _clone_slide(prs, T_FOCUSING_SEGMENTS)
    shapes = _find_text_shapes(slide)

    # Categorize shapes by position
    pct_shapes = []   # auto shapes with percentages (top ~1735234)
    name_shapes = []  # text boxes with names (top ~4242581)
    tag_shapes = []   # text boxes with taglines (top ~4661986+)
    title_shape = None

    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        text = s.text_frame.text.strip()
        if "FOCUSING" in text.upper() or "DOMINANT" in text.upper():
            title_shape = s
        elif text.endswith("%") and len(text) <= 4:
            pct_shapes.append(s)
        elif s.top > 4600000:
            tag_shapes.append(s)
        elif s.top > 4000000 and len(text) < 30:
            name_shapes.append(s)

    # Sort by left position
    pct_shapes.sort(key=lambda s: s.left)
    name_shapes.sort(key=lambda s: s.left)
    tag_shapes.sort(key=lambda s: s.left)

    for i, seg in enumerate(segments[:5]):
        if i < len(pct_shapes):
            _set_text_preserve_format(pct_shapes[i].text_frame, f"{seg.get('size_pct', '?')}%")
        if i < len(name_shapes):
            _set_text_preserve_format(name_shapes[i].text_frame, seg.get("name", f"Segment {i+1}"))
        if i < len(tag_shapes):
            _set_text_preserve_format(tag_shapes[i].text_frame, _truncate(seg.get("tagline", ""), 70))

    # Clear excess shapes if fewer than 5 segments
    for shapes_list in [pct_shapes, name_shapes, tag_shapes]:
        for i in range(len(segments), len(shapes_list)):
            _set_text_preserve_format(shapes_list[i].text_frame, "")

    return slide


def _build_why_not_segments(prs, deprioritized, brand_name):
    """Build 'WHY NOT PRIORITIZE OTHER SEGMENTS (FOR NOW)' slide.

    Uses T_CONTENT template with structured deprioritization rationale.
    """
    slide = _clone_slide(prs, T_CONTENT)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "WHY NOT PRIORITIZE OTHER SEGMENTS (FOR NOW)")

    bullets = []
    for dep in deprioritized[:3]:
        name = dep.get("name", "Segment")
        reason = dep.get("reason", "Not the right fit for now")
        bullets.append(f"{name} ({dep.get('size_pct', '?')}%): {_truncate(reason, 65)}")

    if len(shapes) >= 2:
        _set_text_preserve_format(shapes[1].text_frame, bullets or ["All segments show potential"])
    if len(shapes) >= 3:
        closing = f"Building long-term authority requires focus. {brand_name} must anchor the brand before expanding."
        _set_text_preserve_format(shapes[2].text_frame, _truncate(closing, 85))

    return slide


def _build_competitive_fares(prs, fares_data, brand_name):
    """Build 'HOW [BRAND] FARES AGAINST THE COMPETITION' slide.

    Shows competitive positioning: what each brand wins on, the compromise forced,
    and the strategic question.
    """
    slide = _clone_slide(prs, T_CONTENT)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame,
            _truncate(f"HOW {brand_name.upper()} FARES AGAINST THE COMPETITION", 55))

    brand_strengths = fares_data.get("brand_strengths", "")
    compromise = fares_data.get("category_compromise", "")
    opportunity = fares_data.get("strategic_opportunity", "")

    bullets = []
    if brand_strengths:
        bullets.append(_truncate(brand_strengths, 85))
    if compromise:
        bullets.append(_truncate(compromise, 85))
    if opportunity:
        bullets.append(_truncate(opportunity, 85))

    if len(shapes) >= 2:
        _set_text_preserve_format(shapes[1].text_frame, bullets or ["Competitive analysis in progress"])
    if len(shapes) >= 3:
        question = fares_data.get("strategic_question",
            f"What would it look like to build a brand that doesn't force that compromise?")
        _set_text_preserve_format(shapes[2].text_frame, _truncate(question, 85))

    return slide


def _build_target_recommendation(prs, target):
    """Clone PRIMARY TARGET slide (slide 76 pattern).

    Shape 0: Title — "PRIMARY TARGET: [SEGMENT NAME]"
    Shape 1: Rationale bullets (4 bullets)
    Shape 2: Image (right half)
    Shape 3: Insight text (bottom)
    """
    slide = _clone_slide(prs, T_TARGET_RECOMMENDATION)
    shapes = _find_text_shapes(slide)

    title = target.get("title", "PRIMARY TARGET")

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, _truncate(title, 55))
    if len(shapes) >= 2:
        bullets = target.get("rationale_bullets", [])
        bullets = [_truncate(b, 85) for b in bullets[:4]]
        _set_text_preserve_format(shapes[1].text_frame, bullets)
    if len(shapes) >= 3:
        _set_text_preserve_format(shapes[2].text_frame, _truncate(target.get("insight", ""), 85))

    return slide


def _build_why_target(prs, target):
    """Clone WHY [SEGMENT] slide (slide 77 pattern).

    Shape 0: Title
    Shape 1: Rationale bullets (left)
    Shape 2: Image (right)
    Shape 3: Insight (bottom)
    """
    slide = _clone_slide(prs, T_WHY_TARGET)
    shapes = _find_text_shapes(slide)

    segment_name = target.get("primary_segment", "THIS SEGMENT")

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, _truncate(f"WHY {segment_name.upper()} IS THE RIGHT FOCUS", 55))
    if len(shapes) >= 2:
        bullets = target.get("rationale_bullets", [])
        bullets = [_truncate(b, 85) for b in bullets[:4]]
        _set_text_preserve_format(shapes[1].text_frame, bullets)
    if len(shapes) >= 3:
        _set_text_preserve_format(shapes[2].text_frame, _truncate(target.get("insight", ""), 85))

    return slide


def _build_enables_slide(prs, target):
    """Clone ENABLES slide (slide 78 pattern).

    Shape 0: Title
    Shape 1: "What This Does Not Decide Yet" (right column)
    Shape 2: "What Targeting [X] Unlocks" (left column)
    Shape 3: Closing insight (bottom)
    """
    slide = _clone_slide(prs, T_ENABLES)
    shapes = _find_text_shapes(slide)

    segment_name = target.get("primary_segment", "this segment")
    enables = target.get("enables", [])
    does_not = target.get("does_not_decide", [])

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "WHAT THIS CHOICE ENABLES (AND DOES NOT)")

    # Left column = enables, right column = does not decide
    enables_text = f"What Targeting {segment_name} Unlocks\n" + "\n".join(
        _truncate(e, 85) for e in enables[:3]
    )
    does_not_text = "What This Does Not Decide Yet\n" + "\n".join(
        _truncate(d, 85) for d in does_not[:3]
    )

    if len(shapes) >= 3:
        # shapes sorted by top,left — left column first
        _set_text_preserve_format(shapes[2].text_frame, enables_text)
        _set_text_preserve_format(shapes[1].text_frame, does_not_text)
    if len(shapes) >= 4:
        _set_text_preserve_format(shapes[3].text_frame, _truncate(target.get("insight", ""), 100))

    return slide


def _build_consumer_summary(prs, summary_text):
    """Clone consumer summary slide (slide 79 — half-text, half-image)."""
    slide = _clone_slide(prs, T_CONSUMER_SUMMARY)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "CONSUMER SUMMARY")
    if len(shapes) >= 2:
        _set_text_preserve_format(shapes[1].text_frame, _truncate(summary_text, 260))

    return slide


def _build_final_summary(prs, summary_data):
    """Clone three-column summary slide (slide 80 pattern).

    Shape 0: Title — "SUMMARY & NEXT STEPS"
    Shape 1-3: Column headers (Consumer, Capabilities, Competition)
    Shape 4-6: Column text paragraphs
    Shape 7: Closing insight (bottom)
    """
    slide = _clone_slide(prs, T_FINAL_SUMMARY)
    shapes = _find_text_shapes(slide)

    if len(shapes) >= 1:
        _set_text_preserve_format(shapes[0].text_frame, "SUMMARY & NEXT STEPS")

    cap_text = summary_data.get("capabilities_column", "")
    comp_text = summary_data.get("competition_column", "")
    cons_text = summary_data.get("consumer_column", "")
    closing = summary_data.get("closing_insight", "")

    # Find column header shapes (short text) and body shapes (long text)
    headers = []
    bodies = []
    for s in shapes[1:]:
        text = s.text_frame.text.strip()
        if len(text) < 20:
            headers.append(s)
        elif len(text) > 20:
            bodies.append(s)

    # Set headers
    header_labels = ["Capabilities", "Competition", "Consumer"]
    for i, label in enumerate(header_labels):
        if i < len(headers):
            _set_text_preserve_format(headers[i].text_frame, label)

    # Set body paragraphs
    column_texts = [cap_text, comp_text, cons_text]
    for i, txt in enumerate(column_texts):
        if i < len(bodies):
            _set_text_preserve_format(bodies[i].text_frame, _truncate(txt, 250))

    # Closing insight (last shape with substantial width)
    closing_shapes = [s for s in shapes if s.width > 7000000 and s.top > 4000000]
    if closing_shapes:
        _set_text_preserve_format(closing_shapes[0].text_frame, _truncate(closing, 120))

    return slide


# ── Chart Slide Builders (Questionnaire Section) ──────────────

def _remove_chart_shapes(slide, clean_region=True):
    """Remove chart shapes and optionally all chart-region elements.

    After cloning, chart objects decompose into extra shapes (axis titles,
    connectors, sub-labels). With clean_region=True, we also remove these
    orphaned elements, keeping only the slide title (top < 600000) and
    base/sample text (top > 6000000).
    """
    for shape in list(slide.shapes):
        if shape.shape_type == 3:  # CHART
            shape._element.getparent().remove(shape._element)
        elif shape.shape_type == 9:  # LINE / CONNECTOR
            shape._element.getparent().remove(shape._element)

    if not clean_region:
        return

    for shape in list(slide.shapes):
        if shape.top < 600000 or shape.top > 6000000:
            continue
        shape._element.getparent().remove(shape._element)


def _insert_chart_image(slide, chart_path: Path, left=None, top=None, width=None, height=None):
    """Insert a chart PNG image onto a slide at specified position.

    If position not given, uses default chart area (centered, below title).
    """
    if not chart_path or not chart_path.exists():
        return
    if left is None:
        left = Emu(348906)
    if top is None:
        top = Emu(1879166)
    if width is None:
        width = Emu(11843094)
    if height is None:
        height = Emu(4114800)
    slide.shapes.add_picture(str(chart_path), left, top, width, height)


def _insert_asset_image(slide, asset_name: str, left, top, width, height):
    """Insert an asset image (gender_icon, etc.) at a specific position."""
    asset_path = ASSETS_DIR / asset_name
    if not asset_path.exists():
        return
    slide.shapes.add_picture(str(asset_path), left, top, width, height)


def _build_chart_divider(prs, template_idx, title_override=None):
    """Clone a section divider slide (Demographics, Shopping, Brand Eval)."""
    slide = _clone_slide(prs, template_idx)
    if title_override:
        shapes = _find_text_shapes(slide)
        if shapes:
            _set_text_preserve_format(shapes[0].text_frame, title_override)
    return slide


def _build_chart_slide(prs, chart_data: dict, chart_path: Path, template_idx=None):
    """Build a chart slide by cloning a template and replacing charts with rendered images.

    Args:
        chart_data: Chart metadata from analyzer (title, subtitle, chart_type, etc.)
        chart_path: Path to the rendered chart PNG
        template_idx: Which template slide to clone (auto-selected by chart_type if None)
    """
    chart_type = chart_data.get("chart_type", chart_data.get("type", "hbar"))

    if template_idx is None:
        template_idx = {
            "dual": T_CHART_DUAL,
            "donut": T_CHART_DUAL,
            "pie": T_CHART_SINGLE_HBAR,
            "hbar": T_CHART_SINGLE_HBAR,
            "vbar": T_CHART_SINGLE_VBAR,
            "stacked": T_CHART_STACKED,
            "funnel": T_CHART_SINGLE_HBAR,
            "grouped_bar": T_CHART_SINGLE_HBAR,
            "wordcloud": T_CHART_SINGLE_HBAR,
            "matrix": T_CHART_TABLE,
            "table": T_CHART_TABLE,
        }.get(chart_type, T_CHART_SINGLE_HBAR)

    slide = _clone_slide(prs, template_idx)

    _remove_chart_shapes(slide)

    # Update title (kept by _remove_chart_shapes — top < 600000)
    title = chart_data.get("title", "")
    for s in _find_text_shapes(slide):
        text = s.text_frame.text.strip()
        if text and text.isupper() and len(text) < 80:
            _set_text_preserve_format(s.text_frame, _truncate(title, 60))
            break

    # Add subtitle/question as a new text box (old ones were removed with chart region)
    subtitle = chart_data.get("subtitle", "") or chart_data.get("question", "")
    if subtitle:
        from pptx.util import Pt
        txBox = slide.shapes.add_textbox(
            Emu(419100), Emu(1124334), Emu(11353800), Emu(500000)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = _truncate(subtitle, 120)
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x29, 0x25, 0x24)
        run.font.italic = True

    # Insert rendered chart image below subtitle — size by chart type
    if chart_type == "dual":
        _insert_chart_image(slide, chart_path,
                            left=Emu(419100), top=Emu(1750000),
                            width=Emu(11353800), height=Emu(4300000))
    elif chart_type in ("donut", "pie"):
        # Full-width layout (donut left + bars right) to fill the slide
        _insert_chart_image(slide, chart_path,
                            left=Emu(419100), top=Emu(1750000),
                            width=Emu(11353800), height=Emu(4300000))
    elif chart_type == "vbar":
        # Vertical bar: slightly narrower to avoid edge clipping
        _insert_chart_image(slide, chart_path,
                            left=Emu(600000), top=Emu(1750000),
                            width=Emu(10992000), height=Emu(4300000))
    elif chart_type == "matrix":
        # Matrix/table: full width, slightly taller for readability
        _insert_chart_image(slide, chart_path,
                            left=Emu(419100), top=Emu(1650000),
                            width=Emu(11353800), height=Emu(4600000))
    elif chart_type == "wordcloud":
        # Word cloud: centered, large square area
        chart_w = Emu(8229600)   # 9.0 inches
        chart_h = Emu(4389120)   # 4.8 inches
        chart_left = Emu((12192000 - 8229600) // 2)
        _insert_chart_image(slide, chart_path,
                            left=chart_left, top=Emu(1750000),
                            width=chart_w, height=chart_h)
    elif chart_type == "grouped_bar":
        # Grouped bar: full width with small margins
        _insert_chart_image(slide, chart_path,
                            left=Emu(419100), top=Emu(1750000),
                            width=Emu(11353800), height=Emu(4300000))
    else:
        # Default hbar: full width
        _insert_chart_image(slide, chart_path)

    return slide


def _build_respondent_profile(prs, chart_paths: list[Path], demographics: dict = None):
    """Build respondent profile slide with gender icon + generation/ethnicity charts.

    Clones template slide 28 and replaces charts with rendered images.
    Inserts gender_icon.png asset.
    """
    slide = _clone_slide(prs, 28)
    _remove_chart_shapes(slide, clean_region=True)

    # Remove group shapes (original gender icon groups)
    for shape in list(slide.shapes):
        if shape.shape_type == 6:  # GROUP
            sp = shape._element
            sp.getparent().remove(sp)

    # Insert gender icon
    _insert_asset_image(slide, "gender_icon.png",
                        left=Emu(400000), top=Emu(1400000),
                        width=Emu(2200000), height=Emu(2200000))

    # Insert chart images (generation bar chart on right top, ethnicity on right bottom)
    if len(chart_paths) >= 1:
        _insert_chart_image(slide, chart_paths[0],
                            left=Emu(3364800), top=Emu(1500000),
                            width=Emu(8229600), height=Emu(2100000))
    if len(chart_paths) >= 2:
        _insert_chart_image(slide, chart_paths[1],
                            left=Emu(3364800), top=Emu(3800000),
                            width=Emu(8229600), height=Emu(2100000))

    shapes = _find_text_shapes(slide)
    for s in shapes:
        text = s.text_frame.text.strip()
        if "RESPONDENT" in text.upper():
            _set_text_preserve_format(s.text_frame, "RESPONDENT PROFILE")
            break

    return slide


# ── Main Generator ───────────────────────────────────────────

async def generate_pptx(
    project_id: int,
    analysis: dict,
    brand_name: str,
    phase: str = "full",
    collected_images: dict = None,
) -> tuple[Path, list[dict]]:
    """Generate a Brand Discovery PPTX from analysis data.

    Clones slides from the CozyFit reference template and replaces
    text content with analysis results. If collected_images is provided,
    replaces template images with brand-specific images.

    Args:
        collected_images: Output from image_collector.collect_images()
            {"brand": [Path], "product": [Path], "lifestyle": [Path], "all": [Path]}

    Returns:
        (pptx_path, slide_previews)
    """
    img_pool = _ImagePool(collected_images)

    # Start from the reference PPTX to get its theme, layouts, fonts
    prs = Presentation(str(TEMPLATE_PATH))

    # Remove ALL existing slides — we'll clone fresh ones
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    slide_meta = []
    date_str = analysis.get("date", "2026")

    # ── 1. Cover ──────────────────────────────────────────────
    _build_cover(prs, brand_name, date_str)
    slide_meta.append({"type": "cover", "content": {"brand_name": brand_name}})

    # ── 2. Agenda ─────────────────────────────────────────────
    slide = _build_agenda(prs)
    if img_pool.has_images():
        _replace_slide_image(slide, img_pool.next_brand())
    slide_meta.append({"type": "agenda", "content": {}})

    # ── 3. Approach ───────────────────────────────────────────
    _build_approach(prs)
    slide_meta.append({"type": "approach", "content": {}})

    # ── 4. Step 1 – Discovery ─────────────────────────────────
    _build_step_divider(prs)
    slide_meta.append({"type": "step", "content": {"step": 1}})

    # ── Capabilities ──────────────────────────────────────────

    _build_section_header(prs, "capabilities")
    slide_meta.append({"type": "section", "content": {"section": "capabilities"}})

    cap = analysis.get("capabilities", {})

    # Content slides for each capability dimension
    content_keys = [
        "execution_summary", "product_offer", "product_fundamentals",
        "pricing_position", "channel_analysis",
    ]
    # Alternate between template slide 6 and 7 for visual variety
    template_pool = [T_CONTENT, T_CONTENT_ALT]

    for i, key in enumerate(content_keys):
        section = cap.get(key)
        if section:
            tmpl = template_pool[i % len(template_pool)]
            slide = _build_content_slide(
                prs,
                title=section.get("title", key.replace("_", " ").upper()),
                bullets=section.get("bullets", []),
                insight_text=section.get("insight", ""),
                template_idx=tmpl,
            )
            if img_pool.has_images():
                _replace_slide_image(slide, img_pool.next_brand())
            slide_meta.append({"type": "insight", "content": section})

    # Brand challenges (alternate templates for variety)
    for ci, challenge in enumerate(cap.get("brand_challenges", [])):
        tmpl = template_pool[(ci + 1) % len(template_pool)]
        slide = _build_content_slide(
            prs,
            title=challenge.get("title", "BRAND CHALLENGE"),
            bullets=challenge.get("bullets", []),
            insight_text=challenge.get("insight", ""),
            template_idx=tmpl,
        )
        if img_pool.has_images():
            _replace_slide_image(slide, img_pool.next_brand())
        slide_meta.append({"type": "insight", "content": challenge})

    # Capabilities summary
    cap_summary = cap.get("capabilities_summary", "")
    if cap_summary:
        slide = _build_summary_slide(prs, "CAPABILITIES SUMMARY", cap_summary)
        if img_pool.has_images():
            _replace_slide_image(slide, img_pool.next_lifestyle())
        slide_meta.append({"type": "summary", "content": {"text": cap_summary}})

    # ── Competition (Phase 2+) ────────────────────────────────

    if phase in ("market_structure", "full") and analysis.get("competition"):
        _build_section_header(prs, "competition")
        slide_meta.append({"type": "section", "content": {"section": "competition"}})

        comp = analysis.get("competition", {})

        # Market overview
        overview = comp.get("market_overview", {})
        if overview:
            slide = _build_content_slide(
                prs,
                title=overview.get("title", "COMPETITIVE LANDSCAPE"),
                bullets=overview.get("bullets", []),
                insight_text=overview.get("insight", ""),
            )
            if img_pool.has_images():
                _replace_slide_image(slide, img_pool.next_product())
            slide_meta.append({"type": "insight", "content": overview})

        # Competitor deep dives
        for competitor in comp.get("competitor_analyses", []):
            pos_bullets = [
                f"{p['label']}: {p['detail']}"
                for p in competitor.get("positioning", [])
            ]
            learn_bullets = [
                f"{k['label']}: {k['detail']}"
                for k in competitor.get("key_learnings", [])
            ]
            slide = _build_competitor_slide(
                prs,
                name=competitor.get("name", "Competitor"),
                positioning_bullets=pos_bullets,
                learnings_bullets=learn_bullets,
            )
            if img_pool.has_images():
                _replace_slide_image(slide, img_pool.next_product())
            slide_meta.append({"type": "competitor", "content": competitor})

        # Landscape summary
        landscape = comp.get("landscape_summary", {})
        if landscape:
            roles = landscape.get("market_roles", [])
            role_bullets = [
                f"{r['role']}: {', '.join(r.get('brands', []))} — {r.get('description', '')}"
                for r in roles[:4]
            ]
            white_space = landscape.get("white_space", "")
            sidebar = f"White Space Opportunity:\n{white_space}" if white_space else ""
            slide = _build_landscape_slide(
                prs,
                title="COMPETITIVE LANDSCAPE ROLES",
                bullets=role_bullets or ["No market roles identified"],
                sidebar_text=sidebar,
            )
            slide_meta.append({"type": "landscape", "content": landscape})

        # Competition summary
        comp_summary = comp.get("competition_summary", "")
        if comp_summary:
            slide = _build_summary_slide(prs, "COMPETITION SUMMARY", comp_summary, T_COMP_SUMMARY)
            if img_pool.has_images():
                _replace_slide_image(slide, img_pool.next_lifestyle())
            slide_meta.append({"type": "summary", "content": {"text": comp_summary}})

    # ── Consumer (Full only) ──────────────────────────────────

    if phase == "full" and analysis.get("consumer"):
        _build_section_header(prs, "consumer")
        slide_meta.append({"type": "section", "content": {"section": "consumer"}})

        consumer = analysis.get("consumer", {})

        # Research approach
        research = consumer.get("research_approach", [])
        if research:
            _build_research_approach(prs, research)
            slide_meta.append({"type": "research", "content": {"items": research}})

        # Key consumer insights as content slides
        for insight in consumer.get("key_insights", []):
            slide = _build_content_slide(
                prs,
                title=insight.get("title", "CONSUMER INSIGHT"),
                bullets=insight.get("bullets", []),
                insight_text=insight.get("insight", ""),
            )
            if img_pool.has_images():
                _replace_slide_image(slide, img_pool.next_lifestyle())
            slide_meta.append({"type": "insight", "content": insight})

        # ── Questionnaire / Chart Slides ─────────────────────
        charts = consumer.get("charts", [])
        if charts:
            from pipeline.chart_renderer import render_chart
            chart_output_dir = OUTPUT_DIR / f"project_{project_id}" / "charts"

            # Demographics divider
            _build_chart_divider(prs, T_CHART_DIVIDER_DEMO, "Demographics &\nBackground")
            slide_meta.append({"type": "divider", "content": {"title": "Demographics & Background"}})

            shopping_inserted = False
            drivers_inserted = False
            brand_inserted = False
            metrics_def_inserted = False
            for ci, chart_data in enumerate(charts):
                chart_path = render_chart(chart_data, chart_output_dir, ci)
                chart_type = chart_data.get("chart_type", chart_data.get("type", "hbar"))
                title_lower = chart_data.get("title", "").lower()
                section_lower = chart_data.get("section", "").lower()

                is_shopping = section_lower in ("shopping habits", "shopping") or any(kw in title_lower for kw in ("shopping", "habit", "frequency", "spend", "channel", "purchase channel", "occasion", "pre-purchase", "usage"))
                is_driver = section_lower in ("purchase drivers", "drivers") or any(kw in title_lower for kw in ("driver", "matters most", "premium", "willingness", "pay", "wordcloud", "say about", "pain point"))
                is_brand = section_lower in ("brand evaluation", "brand") or any(kw in title_lower for kw in ("brand", "awareness", "metric", "association", "matrix", "perception", "favorite", "likelihood", "switching"))

                if not shopping_inserted and ci >= 1 and (is_shopping or ci == 4):
                    _build_chart_divider(prs, T_CHART_DIVIDER_SHOPPING, "Shopping Habits, Usage,\nAttitude and Image")
                    slide_meta.append({"type": "divider", "content": {"title": "Shopping Habits"}})
                    shopping_inserted = True
                elif not drivers_inserted and shopping_inserted and (is_driver or ci >= 8):
                    _build_chart_divider(prs, T_CHART_DIVIDER_SHOPPING, "Purchase Drivers\n& Barriers")
                    slide_meta.append({"type": "divider", "content": {"title": "Purchase Drivers"}})
                    drivers_inserted = True
                elif not brand_inserted and (shopping_inserted or drivers_inserted) and (is_brand or ci >= len(charts) - 3):
                    _build_chart_divider(prs, T_CHART_DIVIDER_BRAND, "Brand Evaluation &\nCompetitor Analysis")
                    slide_meta.append({"type": "divider", "content": {"title": "Brand Evaluation"}})
                    brand_inserted = True

                if chart_path is None:
                    continue

                if not metrics_def_inserted and chart_type in ("funnel", "grouped_bar") and is_brand:
                    _build_brand_metrics_def(prs)
                    slide_meta.append({"type": "boilerplate", "content": {"title": "Brand Metrics Definitions"}})
                    metrics_def_inserted = True

                _build_chart_slide(prs, chart_data, chart_path)
                slide_meta.append({"type": "chart", "content": chart_data})

        # Segmentation divider + intro boilerplate
        segments = consumer.get("segments", [])
        if segments:
            _clone_slide(prs, T_SEGMENT_DIVIDER)
            slide_meta.append({"type": "divider", "content": {"title": "Market Segmentation"}})

            _build_segmentation_intro(prs)
            slide_meta.append({"type": "boilerplate", "content": {"title": "Benefits of Segmentation"}})

            # Segment overview (all segments at a glance)
            _build_segment_overview(prs, segments)
            slide_meta.append({"type": "segment_overview", "content": {"segments": [s.get("name") for s in segments]}})

            # "FOCUSING ON THE MOST DOMINANT SEGMENTS…"
            _build_focusing_segments(prs, segments)
            slide_meta.append({"type": "focusing", "content": {"segments": [s.get("name") for s in segments]}})

            # Individual segment slides: 6-slide pattern per segment
            # 1. Meet Segment  2. Respondent Profile  3. Closer Look 1
            # 4. Closer Look 2  5. Challenges Table  6. Closer Look 3
            for seg in segments[:5]:
                slide = _build_meet_segment(prs, seg)
                if img_pool.has_images():
                    _replace_slide_image(slide, img_pool.next_lifestyle(), replace_background=True)
                slide_meta.append({"type": "meet_segment", "content": seg})

                _build_segment_profile(prs, seg)
                slide_meta.append({"type": "segment_profile", "content": {"segment": seg.get("name")}})

                _build_segment_closer_look(prs, seg, slide_num=1)
                slide_meta.append({"type": "closer_look", "content": {"segment": seg.get("name"), "slide": 1}})

                _build_segment_closer_look(prs, seg, slide_num=2)
                slide_meta.append({"type": "closer_look", "content": {"segment": seg.get("name"), "slide": 2}})

                _build_segment_challenges(prs, seg)
                slide_meta.append({"type": "challenges", "content": {"segment": seg.get("name")}})

                slide = _build_segment_closer_look(prs, seg, slide_num=3)
                # Replace all 4 lifestyle card images on closer_look_3
                if img_pool.has_images():
                    _replace_card_images(slide, img_pool)
                slide_meta.append({"type": "closer_look", "content": {"segment": seg.get("name"), "slide": 3}})

        # Target recommendation
        target = consumer.get("target_recommendation", {})
        if target:
            slide = _build_target_recommendation(prs, target)
            if img_pool.has_images():
                _replace_slide_image(slide, img_pool.next_lifestyle())
            slide_meta.append({"type": "target", "content": target})

            _build_why_target(prs, target)
            slide_meta.append({"type": "why_target", "content": target})

            _build_enables_slide(prs, target)
            slide_meta.append({"type": "enables", "content": target})

            # "WHY NOT PRIORITIZE OTHER SEGMENTS"
            deprioritized = consumer.get("deprioritized_segments", [])
            if not deprioritized and len(segments) > 1:
                primary = target.get("primary_segment", "")
                deprioritized = [
                    {"name": s.get("name", ""), "size_pct": s.get("size_pct", "?"),
                     "reason": f"Not the primary focus — different needs and priorities"}
                    for s in segments if s.get("name") != primary
                ][:3]
            if deprioritized:
                _build_why_not_segments(prs, deprioritized, brand_name)
                slide_meta.append({"type": "why_not", "content": {"segments": deprioritized}})

            # "HOW [BRAND] FARES AGAINST THE COMPETITION"
            fares = consumer.get("competitive_fares", {})
            if fares:
                _build_competitive_fares(prs, fares, brand_name)
                slide_meta.append({"type": "fares", "content": fares})

        # Consumer summary
        cons_summary = consumer.get("consumer_summary", "")
        if cons_summary:
            slide = _build_consumer_summary(prs, cons_summary)
            if img_pool.has_images():
                _replace_slide_image(slide, img_pool.next_lifestyle())
            slide_meta.append({"type": "consumer_summary", "content": {"text": cons_summary}})

    # ── Final Summary & Next Steps ───────────────────────────

    summary_data = analysis.get("summary_and_next_steps", {})
    if summary_data and phase == "full":
        _build_final_summary(prs, summary_data)
        slide_meta.append({"type": "final_summary", "content": summary_data})

    # ── Thank You ─────────────────────────────────────────────

    _build_thank_you(prs)
    slide_meta.append({"type": "thank_you", "content": {}})

    # ── Save ──────────────────────────────────────────────────

    # Fix CJK fonts across all slides before saving
    _fix_cjk_fonts(prs)

    output_path = OUTPUT_DIR / f"project_{project_id}" / f"{brand_name}_Brand_Discovery.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    # Re-save to fix XML structure issues from slide cloning (LibreOffice compat)
    prs2 = Presentation(str(output_path))
    prs2.save(str(output_path))

    # Generate previews
    preview_paths = _generate_previews(output_path, project_id)
    for i, pp in enumerate(preview_paths):
        if i < len(slide_meta):
            slide_meta[i]["preview_path"] = str(pp)

    return output_path, slide_meta


# ── Preview Generation ───────────────────────────────────────

def _generate_previews(pptx_path: Path, project_id: int) -> list[Path]:
    """Convert PPTX slides to PNG previews via LibreOffice + PyMuPDF."""
    import subprocess
    import tempfile

    preview_dir = PREVIEW_DIR / f"project_{project_id}"
    preview_dir.mkdir(parents=True, exist_ok=True)

    # Clean old previews
    for old in preview_dir.glob("*.png"):
        old.unlink()

    # Step 1: PPTX -> PDF via LibreOffice
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            subprocess.run([
                "soffice", "--headless", "--convert-to", "pdf",
                "--outdir", tmpdir, str(pptx_path)
            ], capture_output=True, timeout=120)

            pdf_files = list(Path(tmpdir).glob("*.pdf"))
            if not pdf_files:
                return _generate_placeholder_previews(pptx_path, preview_dir)

            # Step 2: PDF -> per-page PNG via PyMuPDF
            import fitz
            doc = fitz.open(str(pdf_files[0]))
            paths = []
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                png_path = preview_dir / f"slide_{i:03d}.png"
                pix.save(str(png_path))
                paths.append(png_path)
            doc.close()
            return paths

    except (FileNotFoundError, subprocess.TimeoutExpired, ImportError):
        pass

    return _generate_placeholder_previews(pptx_path, preview_dir)


def _generate_placeholder_previews(pptx_path: Path, preview_dir: Path) -> list[Path]:
    """Simple placeholder previews when LibreOffice unavailable."""
    from PIL import Image, ImageDraw

    prs = Presentation(str(pptx_path))
    paths = []

    for i, slide in enumerate(prs.slides):
        img = Image.new("RGB", (1280, 720), "#FAFAF9")
        draw = ImageDraw.Draw(img)
        draw.text((20, 20), f"Slide {i + 1}", fill="#E8652D")

        y = 80
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text[:100]
                if text.strip():
                    draw.text((40, y), text, fill="#292524")
                    y += 30
                    if y > 650:
                        break

        path = preview_dir / f"slide_{i:03d}.png"
        img.save(str(path))
        paths.append(path)

    return paths
