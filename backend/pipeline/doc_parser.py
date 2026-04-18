"""Document parsing module.

Parses uploaded PDF, DOCX, PPTX files into structured text
for AI analysis.
"""
from pathlib import Path


async def parse_documents(file_paths: list[str]) -> list[dict]:
    """Parse uploaded documents and extract text content.

    Returns:
        [{"filename": str, "text": str, "tables": [...], "images": [...]}]
    """
    results = []

    for fp in file_paths:
        path = Path(fp)
        if not path.exists():
            continue

        suffix = path.suffix.lower()
        text = ""
        tables = []

        try:
            if suffix == ".pdf":
                text, tables = _parse_pdf(path)
            elif suffix in (".docx", ".doc"):
                text = _parse_docx(path)
            elif suffix == ".txt":
                text = path.read_text(encoding="utf-8", errors="ignore")
            elif suffix == ".pptx":
                text = _parse_pptx(path)
            else:
                text = f"[Unsupported file type: {suffix}]"
        except Exception as e:
            text = f"[Parse error: {e}]"

        results.append({
            "filename": path.name,
            "text": text[:10000],
            "tables": tables,
            "images": [],
        })

    return results


def _parse_pdf(path: Path) -> tuple[str, list]:
    """Extract text and tables from PDF using pdfplumber or PyPDF2."""
    try:
        import pdfplumber
        text_parts = []
        tables = []
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages[:50]:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
                for tbl in (page.extract_tables() or []):
                    if tbl and len(tbl) > 1:
                        headers = [str(cell or "").strip() for cell in tbl[0]]
                        rows = [[str(cell or "").strip() for cell in row] for row in tbl[1:]]
                        tables.append({"headers": headers, "rows": rows})
        return "\n\n".join(text_parts), tables
    except ImportError:
        pass

    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        text = "\n\n".join(
            page.extract_text() or "" for page in reader.pages[:50]
        )
        return text, []
    except ImportError:
        return "[Install pdfplumber or PyPDF2 to parse PDFs]", []


def _parse_docx(path: Path) -> str:
    """Extract text from DOCX."""
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        return "[Install python-docx to parse DOCX files]"


def _parse_pptx(path: Path) -> str:
    """Extract text from PPTX slides."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        texts.append(t)
            if texts:
                parts.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except ImportError:
        return "[Install python-pptx to parse PPTX files]"
